import streamlit as st
import nltk
import string
import torch
import fitz
from docx import Document
from pptx import Presentation
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from sklearn.feature_extraction.text import TfidfVectorizer
from transformers import pipeline, AutoTokenizer


# File Extraction Functions
def extract_text_from_pdf(file): # text extraction from pdf files , for loop to go through each page
    text = ""
    with fitz.open(stream=file.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(file): # text extraction from word files , collects paragraphs joins line by line
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file): # text extraction from powerpoint files , goes through slides and shapes
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# Page setup
st.set_page_config(page_title="Text to Insight", layout="centered")
st.title("🧠 Text to Insight Web App")
st.write("Paste your text below for AI-powered analysis — keywords, sentiment, and summaries.")


# accept user input
user_text = st.text_area("Enter your text here", height=200)

uploaded_file = st.file_uploader(
    "📂 Upload a file (PDF, Word, or PowerPoint)",
    type=["pdf", "docx", "pptx"]
)

# process uploaded file
if uploaded_file: 
    file_type = uploaded_file.name.split('.')[-1].lower() # get file extension
    if file_type == "pdf":
        user_text = extract_text_from_pdf(uploaded_file)
    elif file_type == "docx":
        user_text = extract_text_from_docx(uploaded_file)
    elif file_type == "pptx":
        user_text = extract_text_from_pptx(uploaded_file)
    else:
        st.error("Unsupported file type.")

    with st.expander("👁️ Preview extracted text"):
        st.write(user_text[:10000] + ("..." if len(user_text) > 10000 else ""))


# load and cache resources
@st.cache_resource
def load_nltk_resources():
    nltk.download("stopwords", quiet=True)
    nltk.download("punkt", quiet=True)
    return set(stopwords.words("english")), PorterStemmer()

stop_words, stemmer = load_nltk_resources()


@st.cache_resource
def load_sentiment_model():
    return pipeline("sentiment-analysis", device=0 if torch.cuda.is_available() else -1)


@st.cache_resource
def load_summarization_tools():
    model_name = "sshleifer/distilbart-cnn-12-6"
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    summarizer = pipeline(
        "summarization",
        model=model_name,
        tokenizer=tokenizer,
        device=0 if torch.cuda.is_available() else -1
    )
    return summarizer, tokenizer


# Text preprocessing
def preprocess_text(text):
    text = text.lower()
    text = text.translate(str.maketrans('', '', string.punctuation))
    words = [stemmer.stem(w) for w in text.split() if w not in stop_words]
    return words


# sentiment analysis, keyword extraction, summarization
if st.button("Analyze"):
    if user_text.strip() == "":
        st.warning("Please enter some text or upload a file to analyze.")
    else:
        with st.spinner("Analyzing your text... please wait ⏳"):
            try:
                # --- Preprocess ---
                tokens = preprocess_text(user_text)
                st.subheader("Preprocessed Tokens")
                st.write(tokens[:20], "..." if len(tokens) > 20 else "")

                # --- Keyword Extraction ---
                documents = [user_text.strip()]
                if any(c.isalpha() for c in documents[0]):
                    keyword_vectorizer = TfidfVectorizer(stop_words='english', ngram_range=(1, 2))
                    X = keyword_vectorizer.fit_transform(documents)
                    feature_names = keyword_vectorizer.get_feature_names_out()
                    tfidf_scores = X.toarray()[0]
                    keywords = sorted(zip(feature_names, tfidf_scores), key=lambda x: x[1], reverse=True)[:10]
                    st.subheader("Top Keywords")
                    st.write([kw[0] for kw in keywords])
                else:
                    st.warning("The input text is empty or not valid for keyword extraction.")

                # --- Sentiment Analysis ---
                st.subheader("Sentiment Analysis")
                sentiment_analyzer = load_sentiment_model()
                sentiment_result = sentiment_analyzer(user_text[:512])[0]  # limit to 512 tokens
                label = sentiment_result['label']
                score = sentiment_result['score']
                st.write(f"**Sentiment:** {label} ({score:.3f})")

                # --- Summarization ---
                st.subheader("Summary")
                summarizer, tokenizer = load_summarization_tools()

                max_token_length = 512
                inputs = tokenizer(user_text, return_tensors="pt", truncation=False)
                input_ids = inputs["input_ids"][0]

                chunks = []
                for i in range(0, len(input_ids), max_token_length - 50):  # overlap chunks
                    chunk = input_ids[i:i + max_token_length]
                    text_chunk = tokenizer.decode(chunk, skip_special_tokens=True)
                    chunks.append(text_chunk)

                summaries = []
                for c in chunks:
                    # re-tokenize
                    inputs_chunk = tokenizer(c, return_tensors="pt", truncation=True, max_length=max_token_length)
                    summary = summarizer(
                        tokenizer.decode(inputs_chunk["input_ids"][0], skip_special_tokens=True),
                        max_length=120,
                        min_length=40,
                        do_sample=False
                    )[0]["summary_text"]
                    summaries.append(summary)

                summary_text = " ".join(summaries)
                st.write(summary_text)

            except Exception as e:
                st.error(f"⚠️ Something went wrong: {str(e)}")
